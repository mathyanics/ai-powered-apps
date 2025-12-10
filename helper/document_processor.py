"""
Document Processor for PDF and PowerPoint files
"""

import os
import logging
from typing import List, Dict
import PyPDF2
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles document processing for PDF and PPTX files"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                logger.info(f"Extracting text from PDF with {len(pdf_reader.pages)} pages")
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
                    
            logger.info(f"Extracted {len(text)} characters from PDF")
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise
    
    def extract_text_from_ppt(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            text = ""
            presentation = Presentation(file_path)
            logger.info(f"Extracting text from PPTX with {len(presentation.slides)} slides")
            
            for slide_num, slide in enumerate(presentation.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            text += row_text + "\n"
            
            logger.info(f"Extracted {len(text)} characters from PPTX")
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            raise
    
    def process_document(self, file_path: str, filename: str) -> Dict:
        """
        Process document and return text and chunks
        
        Returns:
            dict with 'text', 'chunks', and 'metadata'
        """
        file_ext = os.path.splitext(filename)[1].lower()
        
        # Extract text based on file type
        if file_ext == '.pdf':
            text = self.extract_text_from_pdf(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            text = self.extract_text_from_ppt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        if not text:
            raise ValueError("No text could be extracted from the document")
        
        # Create document and split into chunks
        doc = Document(
            page_content=text,
            metadata={'source': filename}
        )
        
        chunks = self.text_splitter.split_documents([doc])
        logger.info(f"Document split into {len(chunks)} chunks")
        
        # Calculate metadata
        word_count = len(text.split())
        char_count = len(text)
        reading_time = max(1, word_count // 200)  # ~200 words per minute
        
        return {
            'text': text,
            'chunks': chunks,
            'metadata': {
                'filename': filename,
                'word_count': word_count,
                'character_count': char_count,
                'estimated_reading_time': reading_time,
                'chunk_count': len(chunks)
            }
        }
